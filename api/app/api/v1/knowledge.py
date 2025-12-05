"""
Knowledge Center router - manage knowledge items (videos, PDFs, presentations)
"""
import uuid
import os
import re
import logging
from datetime import datetime
from typing import Optional
from fastapi import APIRouter, Depends, HTTPException, status, UploadFile, File, Form
from fastapi.responses import FileResponse
from sqlalchemy.orm import Session
from sqlalchemy import and_

from app.database import get_db
from app.models import KnowledgeItem, KnowledgeItemType, User, UserRole, IndexUser, Index
from app.schemas.knowledge_item import (
    KnowledgeItemCreate,
    KnowledgeItemUpdate,
    KnowledgeItemResponse,
    KnowledgeItemListResponse
)
from app.api.dependencies import get_current_active_user

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/knowledge", tags=["knowledge"])

# Upload directory for knowledge center files
KNOWLEDGE_UPLOAD_DIR = "/app/uploads/knowledge"
THUMBNAIL_DIR = "/app/uploads/knowledge/thumbnails"


def generate_pdf_thumbnail(pdf_path: str, output_path: str) -> bool:
    """Generate thumbnail from first page of PDF"""
    try:
        from pdf2image import convert_from_path
        from PIL import Image

        # Convert first page to image
        images = convert_from_path(pdf_path, first_page=1, last_page=1, dpi=150)
        if images:
            # Resize to reasonable thumbnail size (max 800px width)
            img = images[0]
            max_width = 800
            if img.width > max_width:
                ratio = max_width / img.width
                new_height = int(img.height * ratio)
                img = img.resize((max_width, new_height), Image.Resampling.LANCZOS)

            # Save as JPEG
            img.save(output_path, 'JPEG', quality=85)
            logger.info(f"Generated PDF thumbnail: {output_path}")
            return True
    except Exception as e:
        logger.error(f"Failed to generate PDF thumbnail: {e}")
    return False


def generate_pptx_thumbnail(pptx_path: str, output_path: str) -> bool:
    """Generate thumbnail from first slide of PPTX"""
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from PIL import Image
        import io

        prs = Presentation(pptx_path)
        if len(prs.slides) > 0:
            slide = prs.slides[0]

            # Check if slide has any images we can use as thumbnail
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    image = shape.image
                    image_bytes = image.blob
                    img = Image.open(io.BytesIO(image_bytes))

                    # Resize to reasonable thumbnail size
                    max_width = 800
                    if img.width > max_width:
                        ratio = max_width / img.width
                        new_height = int(img.height * ratio)
                        img = img.resize((max_width, new_height), Image.Resampling.LANCZOS)

                    # Convert to RGB if necessary
                    if img.mode in ('RGBA', 'P'):
                        img = img.convert('RGB')

                    img.save(output_path, 'JPEG', quality=85)
                    logger.info(f"Generated PPTX thumbnail from image: {output_path}")
                    return True

            # If no images, create a simple placeholder with slide title
            # Create a simple colored rectangle as placeholder
            from PIL import ImageDraw, ImageFont

            img = Image.new('RGB', (800, 450), color=(255, 102, 0))  # Orange background
            draw = ImageDraw.Draw(img)

            # Try to get slide title
            title_text = "PowerPoint"
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        title_text = text[:50] + ('...' if len(text) > 50 else '')
                        break

            # Draw text (simple, no font file needed)
            try:
                # Try to use a default font
                font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 36)
            except:
                font = ImageFont.load_default()

            # Center the text
            bbox = draw.textbbox((0, 0), title_text, font=font)
            text_width = bbox[2] - bbox[0]
            text_height = bbox[3] - bbox[1]
            x = (800 - text_width) // 2
            y = (450 - text_height) // 2
            draw.text((x, y), title_text, fill='white', font=font)

            img.save(output_path, 'JPEG', quality=85)
            logger.info(f"Generated PPTX placeholder thumbnail: {output_path}")
            return True

    except Exception as e:
        logger.error(f"Failed to generate PPTX thumbnail: {e}")
    return False


def can_manage_knowledge(user: User, index_id: str, db: Session) -> bool:
    """Check if user can add/edit/delete knowledge items"""
    # System admins can manage
    if user.role == UserRole.ADMIN:
        return True

    # Index owners can manage
    index_role = db.query(IndexUser).filter(
        and_(
            IndexUser.user_id == user.id,
            IndexUser.index_id == index_id,
            IndexUser.role == 'OWNER'
        )
    ).first()

    return index_role is not None


def can_view_knowledge(user: User, index_id: str, db: Session) -> bool:
    """Check if user can view knowledge items"""
    # System admins can view
    if user.role == UserRole.ADMIN:
        return True

    # Any user with access to the index can view
    index_role = db.query(IndexUser).filter(
        and_(
            IndexUser.user_id == user.id,
            IndexUser.index_id == index_id
        )
    ).first()

    return index_role is not None


def extract_youtube_video_id(url: str) -> Optional[str]:
    """Extract YouTube video ID from various URL formats"""
    patterns = [
        r'(?:youtube\.com\/watch\?v=|youtu\.be\/|youtube\.com\/embed\/)([a-zA-Z0-9_-]{11})',
        r'youtube\.com\/v\/([a-zA-Z0-9_-]{11})',
    ]
    for pattern in patterns:
        match = re.search(pattern, url)
        if match:
            return match.group(1)
    return None


def get_youtube_thumbnail(video_id: str) -> str:
    """Get YouTube video thumbnail URL"""
    return f"https://img.youtube.com/vi/{video_id}/maxresdefault.jpg"


def enrich_knowledge_response(item: KnowledgeItem, db: Session) -> dict:
    """Enrich knowledge item with creator information"""
    item_dict = {
        "id": item.id,
        "title": item.title,
        "description": item.description,
        "content_type": item.content_type.value,
        "content_url": item.content_url,
        "thumbnail_path": item.thumbnail_path,
        "file_name": item.file_name,
        "file_size": item.file_size,
        "index_id": item.index_id,
        "created_by": item.created_by,
        "created_at": item.created_at,
        "updated_at": item.updated_at,
        "display_order": item.display_order,
    }

    # Add creator info
    creator = db.query(User).filter(User.id == item.created_by).first()
    if creator:
        item_dict["creator_name"] = creator.full_name_ar
        item_dict["creator_name_en"] = creator.full_name_en

    return item_dict


@router.get("", response_model=KnowledgeItemListResponse)
async def list_knowledge_items(
    index_id: str,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    List all knowledge items for an index
    Only available for ETARI index type
    """
    # Verify index exists and is ETARI type
    index = db.query(Index).filter(Index.id == index_id).first()
    if not index:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Index not found"
        )

    if index.index_type != 'ETARI':
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Knowledge Center is only available for ETARI indexes"
        )

    # Check permissions
    if not can_view_knowledge(current_user, index_id, db):
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have permission to view knowledge items for this index"
        )

    # Get all items for this index
    items = db.query(KnowledgeItem).filter(
        KnowledgeItem.index_id == index_id
    ).order_by(KnowledgeItem.display_order, KnowledgeItem.created_at.desc()).all()

    enriched_items = [enrich_knowledge_response(item, db) for item in items]

    return KnowledgeItemListResponse(
        items=enriched_items,
        total=len(enriched_items)
    )


@router.post("", response_model=KnowledgeItemResponse)
async def create_knowledge_item(
    index_id: str = Form(...),
    title: str = Form(...),
    description: Optional[str] = Form(None),
    content_type: str = Form(...),
    content_url: Optional[str] = Form(None),
    display_order: int = Form(0),
    file: Optional[UploadFile] = File(None),
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """
    Create a new knowledge item

    - For YouTube: provide content_url with the video link
    - For PDF/PPTX: upload the file
    """
    # Verify index exists and is ETARI type
    index = db.query(Index).filter(Index.id == index_id).first()
    if not index:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Index not found"
        )

    if index.index_type != 'ETARI':
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Knowledge Center is only available for ETARI indexes"
        )

    # Check permissions
    if not can_manage_knowledge(current_user, index_id, db):
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Only admins and index owners can add knowledge items"
        )

    # Validate content type
    try:
        item_type = KnowledgeItemType(content_type)
    except ValueError:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Invalid content type. Must be one of: youtube, pdf, pptx"
        )

    # Generate item ID
    item_id = f"ki_{uuid.uuid4().hex[:12]}"

    # Process based on content type
    final_content_url = None
    thumbnail_path = None
    file_name = None
    file_size = None

    if item_type == KnowledgeItemType.YOUTUBE:
        if not content_url:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="YouTube URL is required for YouTube content type"
            )

        # Extract video ID and validate
        video_id = extract_youtube_video_id(content_url)
        if not video_id:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Invalid YouTube URL"
            )

        final_content_url = content_url
        thumbnail_path = get_youtube_thumbnail(video_id)

    else:  # PDF or PPTX
        if not file:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"File upload is required for {content_type} content type"
            )

        # Validate file extension
        file_ext = os.path.splitext(file.filename)[1].lower()
        expected_ext = f".{content_type}"

        if file_ext != expected_ext:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"File must be a {content_type.upper()} file"
            )

        # Create upload directory
        upload_dir = os.path.join(KNOWLEDGE_UPLOAD_DIR, index_id)
        os.makedirs(upload_dir, exist_ok=True)

        # Save file
        file_name = file.filename
        saved_filename = f"{item_id}{file_ext}"
        file_path = os.path.join(upload_dir, saved_filename)

        try:
            contents = await file.read()
            with open(file_path, "wb") as f:
                f.write(contents)
            file_size = len(contents)
        except Exception as e:
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"Failed to save file: {str(e)}"
            )

        final_content_url = file_path

        # Generate thumbnail for PDF/PPTX
        thumbnail_dir = os.path.join(THUMBNAIL_DIR, index_id)
        os.makedirs(thumbnail_dir, exist_ok=True)
        thumbnail_filename = f"{item_id}_thumb.jpg"
        thumbnail_file_path = os.path.join(thumbnail_dir, thumbnail_filename)

        if item_type == KnowledgeItemType.PDF:
            if generate_pdf_thumbnail(file_path, thumbnail_file_path):
                thumbnail_path = thumbnail_file_path
        elif item_type == KnowledgeItemType.PPTX:
            if generate_pptx_thumbnail(file_path, thumbnail_file_path):
                thumbnail_path = thumbnail_file_path

    # Create knowledge item
    knowledge_item = KnowledgeItem(
        id=item_id,
        title=title,
        description=description,
        content_type=item_type,
        content_url=final_content_url,
        thumbnail_path=thumbnail_path,
        file_name=file_name,
        file_size=file_size,
        index_id=index_id,
        created_by=current_user.id,
        created_at=datetime.utcnow(),
        updated_at=datetime.utcnow(),
        display_order=display_order
    )

    db.add(knowledge_item)

    try:
        db.commit()
        db.refresh(knowledge_item)
    except Exception as e:
        db.rollback()
        # Clean up file if database operation fails
        if final_content_url and os.path.exists(final_content_url):
            os.remove(final_content_url)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to create knowledge item: {str(e)}"
        )

    return enrich_knowledge_response(knowledge_item, db)


@router.get("/{item_id}", response_model=KnowledgeItemResponse)
async def get_knowledge_item(
    item_id: str,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """Get a knowledge item by ID"""
    item = db.query(KnowledgeItem).filter(KnowledgeItem.id == item_id).first()

    if not item:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Knowledge item not found"
        )

    # Check permissions
    if not can_view_knowledge(current_user, item.index_id, db):
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have permission to view this knowledge item"
        )

    return enrich_knowledge_response(item, db)


@router.patch("/{item_id}", response_model=KnowledgeItemResponse)
async def update_knowledge_item(
    item_id: str,
    item_data: KnowledgeItemUpdate,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """Update a knowledge item"""
    item = db.query(KnowledgeItem).filter(KnowledgeItem.id == item_id).first()

    if not item:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Knowledge item not found"
        )

    # Check permissions
    if not can_manage_knowledge(current_user, item.index_id, db):
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Only admins and index owners can update knowledge items"
        )

    # Update fields
    update_dict = item_data.model_dump(exclude_unset=True)

    for field, value in update_dict.items():
        if field == "content_url" and item.content_type == KnowledgeItemType.YOUTUBE:
            # Validate and update YouTube URL
            video_id = extract_youtube_video_id(value)
            if not video_id:
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail="Invalid YouTube URL"
                )
            item.thumbnail_path = get_youtube_thumbnail(video_id)
        setattr(item, field, value)

    item.updated_at = datetime.utcnow()

    try:
        db.commit()
        db.refresh(item)
    except Exception as e:
        db.rollback()
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to update knowledge item: {str(e)}"
        )

    return enrich_knowledge_response(item, db)


@router.delete("/{item_id}")
async def delete_knowledge_item(
    item_id: str,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """Delete a knowledge item"""
    item = db.query(KnowledgeItem).filter(KnowledgeItem.id == item_id).first()

    if not item:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Knowledge item not found"
        )

    # Check permissions
    if not can_manage_knowledge(current_user, item.index_id, db):
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Only admins and index owners can delete knowledge items"
        )

    # Delete file if it exists
    if item.content_type != KnowledgeItemType.YOUTUBE and item.content_url:
        if os.path.exists(item.content_url):
            try:
                os.remove(item.content_url)
            except Exception:
                pass  # Continue even if file deletion fails

    # Delete thumbnail if it exists
    if item.thumbnail_path and item.content_type != KnowledgeItemType.YOUTUBE:
        if os.path.exists(item.thumbnail_path):
            try:
                os.remove(item.thumbnail_path)
            except Exception:
                pass

    db.delete(item)

    try:
        db.commit()
    except Exception as e:
        db.rollback()
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to delete knowledge item: {str(e)}"
        )

    return {"message": "Knowledge item deleted successfully"}


@router.get("/{item_id}/download")
async def download_knowledge_file(
    item_id: str,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """Download a knowledge item file (PDF/PPTX only)"""
    item = db.query(KnowledgeItem).filter(KnowledgeItem.id == item_id).first()

    if not item:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Knowledge item not found"
        )

    # Check permissions
    if not can_view_knowledge(current_user, item.index_id, db):
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have permission to download this file"
        )

    # Check if it's a file type
    if item.content_type == KnowledgeItemType.YOUTUBE:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="YouTube videos cannot be downloaded"
        )

    # Check if file exists
    if not item.content_url or not os.path.exists(item.content_url):
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="File not found"
        )

    # Determine media type
    media_type = "application/pdf" if item.content_type == KnowledgeItemType.PDF else "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    return FileResponse(
        path=item.content_url,
        filename=item.file_name or f"download.{item.content_type.value}",
        media_type=media_type
    )


@router.get("/{item_id}/thumbnail")
async def get_knowledge_thumbnail(
    item_id: str,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """Get thumbnail image for a knowledge item (PDF/PPTX only)"""
    item = db.query(KnowledgeItem).filter(KnowledgeItem.id == item_id).first()

    if not item:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Knowledge item not found"
        )

    # Check permissions
    if not can_view_knowledge(current_user, item.index_id, db):
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have permission to view this thumbnail"
        )

    # YouTube uses external thumbnail URL
    if item.content_type == KnowledgeItemType.YOUTUBE:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Use the thumbnail_path URL directly for YouTube videos"
        )

    # Check if thumbnail exists
    if not item.thumbnail_path or not os.path.exists(item.thumbnail_path):
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Thumbnail not found"
        )

    return FileResponse(
        path=item.thumbnail_path,
        media_type="image/jpeg"
    )
